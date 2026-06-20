# -*- coding: utf-8 -*-
"""Deck: Investigación de mercado — CONSTRUCCIÓN DE PÁGINAS WEB + propuesta de 3 paquetes (LocalExpertiz)."""
import tempfile, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as RGB

NAVY   = RGB(0x0C, 0x1E, 0x35)
NAVY2  = RGB(0x11, 0x28, 0x44)
BLUE   = RGB(0x1A, 0x8E, 0xE6)
GREEN  = RGB(0x22, 0xC4, 0x68)
WHITE  = RGB(0xFF, 0xFF, 0xFF)
LIGHT  = RGB(0xEF, 0xF6, 0xFD)
GRAY   = RGB(0x4B, 0x64, 0x80)
GOLD   = RGB(0xFF, 0xC3, 0x4D)
PALE   = RGB(0xC8, 0xDE, 0xF0)
HEAD = "Poppins"
BODY = "Open Sans"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s

def rect(s, l, t, w, h, color):
    r = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, 18, GRAY, False, BODY)]
    first = True
    for txt, sz, col, bold, font in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(sz); f.bold = bold; f.name = font; f.color.rgb = col
    return tb

def bullets(s, l, t, w, h, items, sz=17, col=GRAY, gap=10, dot=GREEN):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        label, rest = it if isinstance(it, tuple) else (None, it)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "›  "; r.font.size = Pt(sz); r.font.bold = True
        r.font.color.rgb = dot; r.font.name = HEAD
        if label:
            r2 = p.add_run(); r2.text = label + "  "
            r2.font.size = Pt(sz); r2.font.bold = True; r2.font.color.rgb = (WHITE if col == PALE else NAVY); r2.font.name = HEAD
        r3 = p.add_run(); r3.text = rest
        r3.font.size = Pt(sz); r3.font.color.rgb = col; r3.font.name = BODY
    return tb

def heading(s, title, kicker=None, dark=False):
    rect(s, 0.55, 0.6, 0.5, 0.12, GREEN)
    if kicker:
        text(s, 1.15, 0.48, 11, 0.4, [(kicker.upper(), 13, BLUE, True, HEAD)])
    text(s, 0.55, 0.76, 12.3, 1.0, [(title, 30, WHITE if dark else NAVY, True, HEAD)])

def table(s, l, t, w, rows, colw, fs=11.5, hfs=12, rh=0.42, header=NAVY):
    nrows, ncols = len(rows), len(rows[0])
    g = s.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(rh*nrows)).table
    g.first_row = False; g.horz_banding = False
    for ci, cw in enumerate(colw):
        g.columns[ci].width = Inches(cw)
    for ri in range(nrows):
        g.rows[ri].height = Inches(rh)
        for ci in range(ncols):
            c = g.cell(ri, ci)
            c.margin_left = Inches(0.09); c.margin_right = Inches(0.07)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = header if ri == 0 else (WHITE if ri % 2 else LIGHT)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(rows[ri][ci])
            run.font.size = Pt(hfs if ri == 0 else fs)
            run.font.bold = (ri == 0)
            run.font.name = HEAD if ri == 0 else BODY
            run.font.color.rgb = WHITE if ri == 0 else NAVY
    return g

def note(s, txt, t=6.98):
    text(s, 0.55, t, 12.3, 0.45, [(txt, 10.5, GRAY, False, BODY)])

def badge(s, l, t, w, label, color=BLUE):
    rect(s, l, t, w, 0.5, color)
    text(s, l, t, w, 0.5, [(label, 14, WHITE, True, HEAD)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ───────────── S1 Portada ─────────────
s = slide(NAVY)
rect(s, 0, 0, 13.333, 0.18, GREEN)
rect(s, 0, 7.32, 13.333, 0.18, BLUE)
text(s, 1.0, 1.35, 11, 0.5, [("LOCALEXPERTIZ  ·  USA + MX", 15, BLUE, True, HEAD)])
text(s, 1.0, 2.1, 11.4, 2.2,
     [("Mercado de Construcción", 44, WHITE, True, HEAD),
      ("de Páginas Web", 44, GREEN, True, HEAD)])
text(s, 1.0, 4.35, 11.4, 1.0, [("Investigación de precios (MX + USA) · Quién ofrece qué y a cuánto · Modelos de cobro · Propuesta de los 3 paquetes con precios", 17, PALE, False, BODY)])
text(s, 1.0, 6.45, 11, 0.5, [("Reporte de investigación · Junio 2026", 13.5, RGB(0x8F,0xA8,0xC0), False, BODY)])

# ───────────── S2 Objetivo / corrección de enfoque ─────────────
s = slide()
heading(s, "Objetivo y alcance", "Encuadre")
bullets(s, 0.7, 1.8, 12, 3, [
    ("El servicio a vender es la CONSTRUCCIÓN DE PÁGINAS WEB.", "Toda la investigación se centra en ese mercado (no en redes sociales)."),
    ("Se investigan los 3 niveles de servicio", "que definimos, con precios reales en México y Estados Unidos."),
    ("Entregable:", "quién ofrece cada nivel y a cuánto, qué incluye, modelos de cobro y una propuesta de 3 paquetes con precios definidos."),
], sz=17, gap=14)
rect(s, 0.7, 5.0, 12, 1.5, LIGHT)
text(s, 0.95, 5.15, 11.6, 1.2, [
    ("Nota de método: ", 14, NAVY, True, HEAD),
    ("investigación basada en >25 fuentes de mercado (agencias, builders, comparativas y páginas de precios reales) de MX y USA. Las fuentes se listan al final.", 13.5, GRAY, False, BODY),
], anchor=MSO_ANCHOR.MIDDLE)

# ───────────── S3 Los 3 niveles ─────────────
s = slide(NAVY)
heading(s, "Los 3 niveles de servicio", "Qué vamos a vender", dark=True)
badge(s, 0.7, 1.85, 3.7, "NIVEL 1", BLUE)
text(s, 0.7, 2.5, 3.7, 2.6, [("Web llave en mano", 17, WHITE, True, HEAD), ("Construir el sitio y entregarlo. Pago único, sin cuota recurrente.", 13.5, PALE, False, BODY)], space=8)
badge(s, 4.8, 1.85, 3.7, "NIVEL 2  ★", GREEN)
text(s, 4.8, 2.5, 3.7, 2.6, [("Web + soporte mensual", 17, WHITE, True, HEAD), ("Construcción + mantenimiento, seguridad, cambios y soporte recurrente.", 13.5, PALE, False, BODY)], space=8)
badge(s, 8.9, 1.85, 3.7, "NIVEL 3", BLUE)
text(s, 8.9, 2.5, 3.7, 2.6, [("Web con CMS editable", 17, WHITE, True, HEAD), ("Sitio drag-and-drop autoadministrable: el cliente edita su web él mismo.", 13.5, PALE, False, BODY)], space=8)
text(s, 0.7, 5.6, 12, 0.6, [("En las siguientes diapositivas: benchmark de precios y proveedores de cada nivel.", 14, GOLD, False, BODY)])

# ───────────── S4 NIVEL 1 — USA ─────────────
s = slide()
heading(s, "Nivel 1 · Pago único — Estados Unidos", "Benchmark USD")
table(s, 0.6, 1.8, 12.15, [
    ["Proveedor / segmento", "Precio (pago único)", "Notas"],
    ["Sitio PyME 5 páginas (promedio)", "$3,000 – $5,000  (≈ $4,500)", "Diseño responsive, SEO básico, CMS, hasta 10 págs."],
    ["Freelancer", "$1,500 – $8,000", "$50 – $150 / hora"],
    ["Agencia (PyME)", "$5,000 – $30,000", "Sube con integraciones y contenido"],
    ["Agencia US (sitio a medida)", "$12,000 – $45,000", "Marca, design system, migración"],
    ["Tarifa por hora (agencia US)", "$100 – $200 / hora", "Modelo por hora cuando el alcance es incierto"],
], colw=[4.4, 3.5, 4.25], rh=0.6)
note(s, "Fuentes: webfx.com, digitalpresent.io, markbrinker.com, northwestregisteredagent.com, bookipi.com")

# ───────────── S5 NIVEL 1 — México ─────────────
s = slide()
heading(s, "Nivel 1 · Pago único — México", "Benchmark MXN")
table(s, 0.6, 1.8, 12.15, [
    ["Tipo de sitio", "Precio (pago único MXN)", "Notas"],
    ["Landing page", "$10,000 – $18,000", "Plantilla, 3–5 secciones, formulario, responsive"],
    ["Sitio corporativo / PyME", "$18,000 – $25,000", "5–10 páginas, blog, SEO básico"],
    ["E-commerce (WooCommerce)", "$25,000 – $50,000+", "Tienda, productos iniciales, pasarelas MX"],
    ["Diseño a medida UX/UI", "desde $30,000", "Sitio personalizado de alto nivel"],
    ["Tarifa frontend por hora", "$500 – $2,000 / hora", "Según complejidad"],
], colw=[4.2, 3.6, 4.35], rh=0.6)
note(s, "Fuentes: mexicowordpress.com, disenador-web-mexico.com, bigredes.com, magokoro.mx, listoweb.com.mx")

# ───────────── S6 NIVEL 2 — qué incluye + USA ─────────────
s = slide()
heading(s, "Nivel 2 · Web + soporte mensual", "Mantenimiento / care plans (USA)")
text(s, 0.6, 1.7, 12, 0.5, [("Incluye típicamente: ", 13.5, NAVY, True, HEAD), ("hosting, actualizaciones, seguridad, backups, monitoreo de uptime, horas de cambios/contenido y soporte.", 13.5, GRAY, False, BODY)])
table(s, 0.6, 2.25, 12.15, [
    ["Proveedor (USA)", "Precio / mes", "Incluye"],
    ["WebyKing — Startup / Growth / Enterprise", "$129 / $349 / $899", "Backups, seguridad, 5–25 h soporte, tareas"],
    ["Integral Web Designs (pay-monthly, sin anticipo)", "$78 / $104 / $198", "Sitio + hosting + dominio + updates + SSL"],
    ["Rango general PyME (care plan)", "$75 – $300 / mes", "Lo común; básico $50, premium hasta $500"],
], colw=[5.0, 2.7, 4.45], rh=0.7)
note(s, "Fuentes: webyking.com, integralwebdesigns.com, webstacks.com, onthemap.com, jim.com")

# ───────────── S7 NIVEL 2 — México ─────────────
s = slide()
heading(s, "Nivel 2 · Web + soporte mensual", "Mantenimiento (México)")
table(s, 0.6, 1.8, 12.15, [
    ["Nivel de mantenimiento", "Precio / mes (MXN)", "Incluye"],
    ["Básico", "$2,500 – $8,000", "Actualizaciones, backups, seguridad, hosting"],
    ["Avanzado", "$10,000 – $35,000", "Soporte, monitoreo 24/7, optimización"],
    ["Rango más común (PyME)", "$1,000 – $5,000", "Lo habitual para sitios de negocio"],
    ["Hosting (referencia)", "desde ~$100 / mes", "Costo base de servidor"],
], colw=[4.3, 3.4, 4.45], rh=0.62)
note(s, "Fuentes: nerade.com, magokoro.mx, cronoshare.com.mx, godaddy.com, mexicowordpress.com")

# ───────────── S8 NIVEL 3 — builders DIY ─────────────
s = slide()
heading(s, "Nivel 3 · CMS editable (drag-and-drop)", "Plataformas / builders — USD/mes")
table(s, 0.6, 1.8, 12.15, [
    ["Plataforma", "Precio / mes", "Para qué sirve / nota"],
    ["Framer", "$5 – $10", "Diseño moderno no-code, fácil de editar"],
    ["Wix", "desde ~$17", "El más amigable para principiantes"],
    ["Squarespace", "$16 / $23 / $39 / $99", "Plantillas + editor Fluid Engine"],
    ["Webflow", "$15 (Basic) / $25 (CMS)", "Máximo control de diseño + CMS visual"],
    ["Webflow E-commerce", "$29 / $74 / $212", "Tienda con CMS visual"],
], colw=[3.6, 3.4, 5.15], rh=0.6)
note(s, "El cliente edita su contenido; la cuota de plataforma corre aparte o gestionada. Fuentes: webflow.com, squarespace.com, wix.com, comparetiers.com")

# ───────────── S9 NIVEL 3 — agencias autoadministrable ─────────────
s = slide()
heading(s, "Nivel 3 · Quién entrega webs autoadministrables", "Agencias (MX + USA)")
table(s, 0.6, 1.8, 12.15, [
    ["Proveedor", "Mercado", "Precio", "Nota"],
    ["Ve mi Página", "MX", "desde $4,999 MXN (pago único)", "WordPress autoadministrable + hosting/dominio"],
    ["Mexico WordPress", "MX", "$10,000 – $50,000 MXN", "WP + Elementor 100% autoeditable"],
    ["Sicom Web", "MX", "cotización", "Panel de administración para editar contenido"],
    ["Agencias WaaS (USA)", "USA", "$149 – $5,000 / mes", "Suscripción: construyen y editas/solicitas"],
], colw=[3.0, 1.3, 3.4, 4.45], rh=0.62)
note(s, "Fuentes: vemipagina.com, mexicowordpress.com, sicomweb.com.mx, rubik.design, penji.co")

# ───────────── S10 Modelos de cobro ─────────────
s = slide(NAVY)
heading(s, "Modelos de cobro en desarrollo web", "Cómo cobra el mercado", dark=True)
bullets(s, 0.7, 1.8, 12, 4, [
    ("Paquete / precio fijo (dominante, 82%):", "alcance cerrado y precio único — ideal para Nivel 1 y 3."),
    ("Suscripción / retainer:", "cuota mensual recurrente — base del Nivel 2 (soporte) y del modelo WaaS."),
    ("Por hora:", "$100–$200 USD/h (USA) · $500–$2,000 MXN/h (MX) — para alcance incierto o extras."),
    ("Pago mensual sin anticipo (pay-monthly / WaaS):", "$0 inicial + cuota — baja la barrera de entrada."),
    ("Value-based:", "precio según el valor/resultados — lo usan los que más facturan."),
], sz=16, col=PALE, gap=12)

# ───────────── S11 Matriz competidores ─────────────
s = slide()
heading(s, "Matriz de competidores (web)", "Panorama")
table(s, 0.55, 1.8, 12.25, [
    ["Competidor", "Mercado", "Oferta", "Precio", "Modelo"],
    ["Ve mi Página", "MX", "Web autoadministrable WP", "desde $4,999 MXN", "Pago único"],
    ["Mexico WordPress", "MX", "WP a medida autoeditable", "$10k–$50k MXN", "Pago único"],
    ["Integral Web Designs", "USA", "Web + soporte (sin anticipo)", "$78–$198 / mes", "Suscripción"],
    ["WebyKing", "USA", "Care plans / mantenimiento", "$129–$899 / mes", "Retainer"],
    ["Webflow / Wix / Squarespace", "Global", "Builder CMS (DIY)", "$5–$99 / mes", "Suscripción"],
    ["Agencias US a medida", "USA", "Sitio personalizado", "$3k–$45k", "Pago único"],
], colw=[3.0, 1.2, 3.4, 2.45, 2.2], rh=0.5, fs=11, hfs=11.5)

# ───────────── S12 Hallazgos / oportunidades ─────────────
s = slide()
heading(s, "Hallazgos y oportunidades", "Dónde ganamos")
bullets(s, 0.7, 1.8, 12, 4, [
    ("MX es mucho más barato que USA:", "mismo trabajo se cobra ~5–10x más en USD → oportunidad de arbitraje bilingüe."),
    ("Pocos ofrecen los 3 niveles juntos:", "empaquetarlos (único / con soporte / autoeditable) nos diferencia."),
    ("El modelo 'sin anticipo + mensual' está creciendo:", "baja la barrera y crea ingreso recurrente."),
    ("La transparencia de precios es ventaja:", "publicar paquetes claros genera confianza y filtra prospectos."),
    ("El soporte recurrente es el verdadero negocio:", "el Nivel 2 da ingresos predecibles mes a mes."),
], sz=16, gap=11)

# ───────────── S13 Paquete 1 ─────────────
s = slide()
heading(s, "Propuesta · Paquete 1 — WEB LANZAMIENTO", "Pago único · llave en mano")
text(s, 0.7, 1.65, 12, 0.5, [("Para: ", 14, NAVY, True, HEAD), ("negocios que necesitan una web profesional y autosuficiente, sin cuota recurrente.", 14, GRAY, False, BODY)])
bullets(s, 0.7, 2.25, 7.4, 3.6, [
    "Diseño a medida responsive (hasta 6 secciones/páginas).",
    "SEO on-page básico + formulario + botón de WhatsApp.",
    "Google Analytics y carga de contenido inicial.",
    "Dominio + hosting gestionado (1er año) y capacitación básica.",
    "2 rondas de revisión. Pago 50% anticipo / 50% entrega.",
], sz=14.5, gap=9)
rect(s, 8.5, 2.25, 4.25, 3.4, LIGHT)
text(s, 8.75, 2.45, 3.8, 3.1, [
    ("PRECIO", 13, BLUE, True, HEAD),
    ("México", 14, NAVY, True, HEAD),
    ("desde $18,000 MXN  ($18,000–$26,000)", 14, GRAY, False, BODY),
    ("Estados Unidos", 14, NAVY, True, HEAD),
    ("desde $2,500 USD  ($2,500–$3,900)", 14, GRAY, False, BODY),
    ("Modelo: pago único (50/50)", 12.5, GREEN, True, HEAD),
], space=10)
note(s, "Posicionamiento: por encima del freelancer económico, por debajo de la agencia premium. Precio 'todo incluido'.")

# ───────────── S14 Paquete 2 ─────────────
s = slide(NAVY)
heading(s, "Propuesta · Paquete 2 — WEB + CRECIMIENTO  ★", "Web + soporte mensual (RECOMENDADO)", dark=True)
text(s, 0.7, 1.65, 12, 0.5, [("Para: ", 14, WHITE, True, HEAD), ("negocios que quieren su web siempre actualizada, segura y con cambios sin preocuparse.", 14, PALE, False, BODY)])
bullets(s, 0.7, 2.25, 7.4, 3.6, [
    "Todo lo del Paquete 1, más plan mensual:",
    "Hosting premium, actualizaciones, seguridad y backups.",
    "Monitoreo de uptime + hasta 2 h de cambios/contenido al mes.",
    "Reporte trimestral y soporte prioritario.",
], sz=14.5, gap=9, col=PALE, dot=GREEN)
rect(s, 8.5, 2.2, 4.25, 4.0, NAVY2)
text(s, 8.75, 2.4, 3.8, 3.7, [
    ("PRECIO", 13, GREEN, True, HEAD),
    ("Opción A — build + plan", 13.5, WHITE, True, HEAD),
    ("MX: desde $20,000 + $3,000–$4,500/mes", 12.5, PALE, False, BODY),
    ("USA: desde $3,000 + $150–$300/mes", 12.5, PALE, False, BODY),
    ("Opción B — sin anticipo (12 m)", 13.5, WHITE, True, HEAD),
    ("MX: $3,800–$5,500 / mes", 12.5, PALE, False, BODY),
    ("USA: $180–$280 / mes", 12.5, PALE, False, BODY),
    ("Modelo: suscripción / retainer", 12, GREEN, True, HEAD),
], space=8)
note(s, "Es el paquete a destacar: ingreso recurrente y relación de largo plazo. Anclado a WebyKing ($129–$899) e Integral Web ($78–$198).")

# ───────────── S15 Paquete 3 ─────────────
s = slide()
heading(s, "Propuesta · Paquete 3 — WEB AUTÓNOMA", "CMS editable drag-and-drop")
text(s, 0.7, 1.65, 12, 0.5, [("Para: ", 14, NAVY, True, HEAD), ("clientes que quieren editar su web ellos mismos, sin depender de nadie.", 14, GRAY, False, BODY)])
bullets(s, 0.7, 2.25, 7.4, 3.6, [
    "Web construida en CMS editable (Webflow / WordPress+Elementor / Framer).",
    "Entregada 100% autoadministrable.",
    "Capacitación + video tutorial + documentación.",
    "Accesos de editor y 30 días de soporte post-entrega.",
], sz=14.5, gap=9)
rect(s, 8.5, 2.25, 4.25, 3.7, LIGHT)
text(s, 8.75, 2.45, 3.8, 3.4, [
    ("PRECIO", 13, BLUE, True, HEAD),
    ("México", 14, NAVY, True, HEAD),
    ("$24,000 – $38,000 MXN", 14, GRAY, False, BODY),
    ("Estados Unidos", 14, NAVY, True, HEAD),
    ("$3,500 – $6,000 USD", 14, GRAY, False, BODY),
    ("Plataforma (aparte/gestionada):", 12, NAVY, True, HEAD),
    ("Webflow $15–$25 · Wix $17 · Squarespace $16–$39 /mes", 11.5, GRAY, False, BODY),
], space=9)
note(s, "Opción de añadir el soporte del Paquete 2. Más caro que el P1 por la capacitación y el setup del CMS editable.")

# ───────────── S16 Resumen de precios ─────────────
s = slide()
heading(s, "Resumen de los 3 paquetes", "Precios propuestos")
table(s, 0.55, 1.85, 12.25, [
    ["Paquete", "Modelo", "Precio México", "Precio USA"],
    ["1 · WEB LANZAMIENTO", "Pago único", "desde $18,000 MXN", "desde $2,500 USD"],
    ["2 · WEB + CRECIMIENTO  ★", "Build + mensual / suscripción", "desde $20,000 + $3,000–$4,500/mes  (o $3,800–$5,500/mes)", "desde $3,000 + $150–$300/mes  (o $180–$280/mes)"],
    ["3 · WEB AUTÓNOMA", "Pago único + plataforma", "$24,000–$38,000 MXN", "$3,500–$6,000 USD"],
], colw=[3.1, 3.0, 3.2, 2.95], rh=0.95, fs=11.5)
text(s, 0.55, 5.75, 12.2, 0.9, [("Recomendación: ", 14, NAVY, True, HEAD), ("liderar la venta con el Paquete 2 (ingreso recurrente), usar el Paquete 1 como entrada y el Paquete 3 como diferenciador para quien quiere autonomía.", 14, GRAY, False, BODY)])

# ───────────── S17 Próximos pasos ─────────────
s = slide(NAVY)
heading(s, "Próximos pasos", "Cierre", dark=True)
bullets(s, 0.7, 1.9, 12, 3.5, [
    ("Validar precios y márgenes", "contra nuestros costos reales (horas, plataforma, hosting)."),
    ("Construir la sección de Paquetes", "en el sitio con estos 3 niveles y precios claros."),
    ("Definir el guion de venta", "que lleve al Paquete 2 (recurrente) como opción recomendada."),
    ("Cargar casos/portafolio web reales", "para respaldar el precio con prueba."),
], sz=16, col=PALE, gap=13)
text(s, 0.7, 6.4, 12, 0.5, [("“Sigamos haciendo que las cosas sucedan.”", 18, GOLD, True, HEAD)])

# ───────────── S18 Fuentes ─────────────
s = slide()
heading(s, "Fuentes consultadas", "Referencias")
text(s, 0.7, 1.8, 6.0, 5, [
    ("Estados Unidos", 14, BLUE, True, HEAD),
    ("webfx.com · digitalpresent.io · webyking.com", 12.5, GRAY, False, BODY),
    ("integralwebdesigns.com · webstacks.com", 12.5, GRAY, False, BODY),
    ("markbrinker.com · northwestregisteredagent.com", 12.5, GRAY, False, BODY),
    ("brandvm.com · penji.co · rubik.design", 12.5, GRAY, False, BODY),
    ("Plataformas / builders", 14, BLUE, True, HEAD),
    ("webflow.com · squarespace.com · wix.com", 12.5, GRAY, False, BODY),
    ("framer (comparetiers.com)", 12.5, GRAY, False, BODY),
], space=8)
text(s, 7.0, 1.8, 5.8, 5, [
    ("México", 14, GREEN, True, HEAD),
    ("mexicowordpress.com · disenador-web-mexico.com", 12.5, GRAY, False, BODY),
    ("bigredes.com · magokoro.mx · listoweb.com.mx", 12.5, GRAY, False, BODY),
    ("vemipagina.com · sicomweb.com.mx · nerade.com", 12.5, GRAY, False, BODY),
    ("cronoshare.com.mx · godaddy.com (LATAM)", 12.5, GRAY, False, BODY),
    ("Nota", 14, NAVY, True, HEAD),
    ("Precios de mercado a jun-2026; varían por alcance,", 12.5, GRAY, False, BODY),
    ("complejidad y proveedor. Rangos, no cotización fija.", 12.5, GRAY, False, BODY),
], space=8)

out = os.path.join(tempfile.gettempdir(), "LocalExpertiz-Investigacion-Web.pptx")
prs.save(out)
print("SAVED:", out)
print("SLIDES:", len(prs.slides._sldIdLst))
print("SIZE_KB:", round(os.path.getsize(out)/1024, 1))
